from __future__ import annotations

import csv
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "reports" / "roadshow_ppt"
ASSET_DIR = ROOT / "reports" / "assets"
IMAGE2_DIR = OUT_DIR / "image2_slides"
OUT_PATH = ROOT / "reports" / "路演PPT.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BLUE = RGBColor(22, 78, 158)
DEEP_BLUE = RGBColor(13, 47, 107)
CYAN = RGBColor(8, 145, 178)
GREEN = RGBColor(5, 150, 105)
AMBER = RGBColor(180, 83, 9)
PURPLE = RGBColor(109, 40, 217)
INK = RGBColor(15, 23, 42)
TEXT = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
BG = RGBColor(248, 251, 255)
PANEL = RGBColor(255, 255, 255)
BORDER = RGBColor(203, 213, 225)


def cm(v: float) -> float:
    return v / 2.54


def image_size(path: Path) -> tuple[int, int]:
    with Image.open(path) as img:
        return img.size


def add_picture_fit(slide, path: Path, x, y, w, h) -> None:
    if not path.exists():
        add_note(slide, f"缺失图片：{path}", x, y, w, h)
        return
    iw, ih = image_size(path)
    box_ratio = float(w) / float(h)
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        width = w
        height = int(w / img_ratio)
        left = x
        top = y + int((h - height) / 2)
    else:
        height = h
        width = int(h * img_ratio)
        left = x + int((w - width) / 2)
        top = y
    slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def add_full_slide_image(slide, path: Path) -> None:
    slide.shapes.add_picture(str(path), 0, 0, width=SLIDE_W, height=SLIDE_H)


def set_fill(shape, color: RGBColor, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def set_line(shape, color: RGBColor = BORDER, width: float = 1.0, transparency: int = 0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 24,
    bold: bool = False,
    color: RGBColor = TEXT,
    align=PP_ALIGN.LEFT,
    font: str = "Microsoft YaHei",
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float, *, size: int = 18, color: RGBColor = TEXT) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)


def add_title(slide, title: str, subtitle: str | None = None, *, page: str | None = None) -> None:
    add_text(slide, title, 0.55, 0.34, 8.8, 0.55, size=28, bold=True, color=DEEP_BLUE)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.95), Inches(1.55), Inches(0.045))
    set_fill(line, BLUE)
    set_line(line, BLUE)
    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.1), Inches(0.95), Inches(0.42), Inches(0.045))
    set_fill(line2, GREEN)
    set_line(line2, GREEN)
    if subtitle:
        add_text(slide, subtitle, 0.55, 1.06, 8.8, 0.28, size=11, color=MUTED)
    if page:
        add_text(slide, page, 12.15, 0.34, 0.75, 0.3, size=11, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)


def add_bg(slide) -> None:
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_fill(rect, BG)
    set_line(rect, BG)
    # Light blueprint grid accents.
    for x in [0.65, 1.95, 3.25, 4.55, 5.85, 7.15, 8.45, 9.75, 11.05, 12.35]:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(1.5), Inches(x), Inches(7.0))
        line.line.color.rgb = RGBColor(230, 238, 247)
        line.line.width = Pt(0.5)
    for y in [1.55, 2.55, 3.55, 4.55, 5.55, 6.55]:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.35), Inches(y), Inches(12.95), Inches(y))
        line.line.color.rgb = RGBColor(230, 238, 247)
        line.line.width = Pt(0.5)


def add_panel(slide, x, y, w, h, *, accent: RGBColor = BLUE, fill: RGBColor = PANEL, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    panel = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(panel, fill)
    set_line(panel, BORDER, 1.0)
    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(h))
    set_fill(marker, accent)
    set_line(marker, accent)
    return panel


def add_note(slide, text, x, y, w, h, *, color=AMBER) -> None:
    add_panel(slide, x, y, w, h, accent=color)
    add_text(slide, text, x + 0.18, y + 0.22, w - 0.36, h - 0.36, size=16, bold=True, color=color)


def add_metric(slide, x, y, w, h, label, value, accent=BLUE) -> None:
    add_panel(slide, x, y, w, h, accent=accent)
    add_text(slide, value, x + 0.18, y + 0.23, w - 0.36, 0.55, size=30, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.15, y + 0.88, w - 0.3, 0.34, size=12, bold=True, color=TEXT, align=PP_ALIGN.CENTER)


def add_card(slide, x, y, w, h, title, body: list[str], accent=BLUE, title_size=19, body_size=15) -> None:
    add_panel(slide, x, y, w, h, accent=accent)
    add_text(slide, title, x + 0.22, y + 0.18, w - 0.36, 0.35, size=title_size, bold=True, color=accent)
    add_bullets(slide, body, x + 0.25, y + 0.68, w - 0.45, h - 0.82, size=body_size)


def arrow(slide, x1, y1, x2, y2, color=CYAN) -> None:
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(2.2)
    conn.line.end_arrowhead = True


def add_flow_step(slide, x, y, w, h, title, subtitle, accent=BLUE) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, RGBColor(255, 255, 255))
    set_line(shape, accent, 1.2)
    add_text(slide, title, x + 0.12, y + 0.18, w - 0.24, 0.3, size=15, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text(slide, subtitle, x + 0.14, y + 0.58, w - 0.28, h - 0.7, size=11, color=TEXT, align=PP_ALIGN.CENTER)


def parse_results() -> tuple[dict[str, float], dict[str, str]]:
    results_path = ROOT / "train" / "results.csv"
    with results_path.open("r", encoding="utf-8", newline="") as f:
        rows = list(csv.DictReader(f))
    final = {k.strip(): float(v) for k, v in rows[-1].items()}
    args: dict[str, str] = {}
    for line in (ROOT / "train" / "args.yaml").read_text(encoding="utf-8").splitlines():
        if ":" in line and not line.lstrip().startswith("#"):
            k, v = line.split(":", 1)
            args[k.strip()] = v.strip()
    return final, args


def slide_full_image(prs: Presentation, path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_slide_image(slide, path)


def build_ppt() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    final, args = parse_results()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Cover generated by image2.
    slide_full_image(prs, IMAGE2_DIR / "cover_b_style_image2.png")

    # 2. Section generated by image2.
    slide_full_image(prs, IMAGE2_DIR / "section_01_background_solution_image2.png")

    # 3. Background and pain points.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "项目背景与痛点", "地下停车场空间复杂，车位分散，摄像头数据需要转化为可查询轨迹", page="03")
    add_card(slide, 0.65, 1.65, 3.8, 4.8, "真实痛点", [
        "车主或管理人员只能依赖记忆、人工巡查或监控回放",
        "多点摄像头各自孤立，缺少以车牌为索引的查询入口",
        "地下环境存在遮挡、反光、低照度和路径复杂等挑战",
    ], BLUE, 21, 16)
    add_card(slide, 4.75, 1.65, 3.8, 4.8, "建设目标", [
        "从视频帧中识别车牌并记录最后出现位置",
        "支持 PC 开发验证和 RDK X5 板端部署",
        "提供 Web 控制台完成监控、日志、查询和路线展示",
    ], GREEN, 21, 16)
    add_picture_fit(slide, ROOT / "docs" / "tutorial-assets" / "garage-ai-system.png", Inches(8.85), Inches(1.65), Inches(3.85), Inches(4.8))

    # 4. Solution overview.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "方案总览：从车牌到位置", "一条清晰的数据链路把“看到车辆”转化为“可查询位置”", page="04")
    steps = [
        ("多路输入", "摄像头 / RTSP\n视频 / 图片", BLUE),
        ("车牌检测", "YOLO Pose\n车牌框 + 四角点", CYAN),
        ("字符识别", "PaddleOCR 或\nLPRNet", GREEN),
        ("轨迹入库", "SQLite 保存\n车牌/摄像头/时间", PURPLE),
        ("Web 查询", "最后位置\n路线可视化", AMBER),
    ]
    x0 = 0.8
    for i, (title, subtitle, accent) in enumerate(steps):
        x = x0 + i * 2.45
        add_flow_step(slide, x, 2.1, 1.85, 1.35, title, subtitle, accent)
        if i < len(steps) - 1:
            arrow(slide, x + 1.86, 2.78, x + 2.38, 2.78, CYAN)
    add_card(slide, 0.8, 4.35, 3.7, 1.65, "核心结果", ["输入车牌号后返回最后出现摄像头、时间和车牌裁切图。"], BLUE, 18, 15)
    add_card(slide, 4.85, 4.35, 3.7, 1.65, "系统形态", ["Python 后端 + SQLite 数据库 + Web 控制台，无需复杂前端构建。"], GREEN, 18, 15)
    add_card(slide, 8.9, 4.35, 3.7, 1.65, "部署路径", ["PC 联调验证，RDK X5 BPU 进行边缘侧现场部署。"], AMBER, 18, 15)

    # 5. Technical section generated by image2.
    slide_full_image(prs, IMAGE2_DIR / "section_02_technical_implementation_image2.png")

    # 6. Architecture.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "系统总体架构", "采集、推理、存储、Web 查询解耦，便于 PC/RDK 双端复用", page="06")
    add_picture_fit(slide, ASSET_DIR / "system_architecture_diagram.png", Inches(0.7), Inches(1.4), Inches(7.7), Inches(5.65))
    add_card(slide, 8.65, 1.55, 3.9, 1.25, "入口层", ["web_app.py 启动 HTTP 服务和推理运行时"], BLUE, 17, 13)
    add_card(slide, 8.65, 3.0, 3.9, 1.25, "后端抽象", ["utils/inference.py 统一 PCDetectionBackend 与 BPUDetectionBackend"], GREEN, 17, 13)
    add_card(slide, 8.65, 4.45, 3.9, 1.25, "数据层", ["DBManager 写入 vehicle_history 并支持精确/模糊查询"], PURPLE, 17, 13)
    add_card(slide, 8.65, 5.9, 3.9, 0.9, "前端层", ["Web 页面展示监控、日志、查询结果和路线"], AMBER, 17, 13)

    # 7. End-to-end logic.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "端到端实现逻辑", "每个环节职责单一：采集只拿帧，推理只识别，接口只读快照或查库", page="07")
    add_picture_fit(slide, ASSET_DIR / "end_to_end_logic_diagram.png", Inches(0.65), Inches(1.5), Inches(7.7), Inches(5.4))
    add_card(slide, 8.7, 1.6, 3.9, 1.2, "最新帧策略", ["每路摄像头只保留最新帧，避免视频队列积压导致定位延迟。"], CYAN, 17, 13)
    add_card(slide, 8.7, 3.05, 3.9, 1.2, "统一结果对象", ["DetectionResult 封装 box、pts、text、confidence、crop。"], BLUE, 17, 13)
    add_card(slide, 8.7, 4.5, 3.9, 1.2, "查询闭环", ["SQLite 记录历史轨迹，/api/search 返回最后出现位置。"], GREEN, 17, 13)
    add_card(slide, 8.7, 5.95, 3.9, 0.9, "可维护性", ["各层可单独测试和替换。"], AMBER, 17, 13)

    # 8. Inference backend.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "推理后端设计", "同一套上层逻辑，兼容 PC 开发验证和 RDK X5 边缘部署", page="08")
    add_picture_fit(slide, ASSET_DIR / "backend_pipeline_diagram.png", Inches(0.65), Inches(1.45), Inches(7.6), Inches(4.1))
    add_card(slide, 8.55, 1.45, 4.1, 1.65, "PC 后端", [
        "YOLO .pt 输出检测框和四角点",
        "PaddleOCR 识别车牌文本",
        "适合快速迭代与页面联调",
    ], BLUE, 18, 14)
    add_card(slide, 8.55, 3.35, 4.1, 1.65, "RDK BPU 后端", [
        "BGR -> NV12 后送入 YOLO .bin",
        "后处理解码 + NMS + LPRNet .bin",
        "面向边缘端低功耗部署",
    ], GREEN, 18, 14)
    add_card(slide, 0.9, 5.85, 11.3, 0.9, "设计收益", [
        "WebInferenceRuntime 只调用 detector.detect(frame)，不关心模型运行平台，降低后续替换模型或迁移硬件的成本。",
    ], AMBER, 17, 14)

    # 9. Threading.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "线程安排与通信机制", "采集线程、推理线程、HTTP 请求线程之间通过缓冲区和运行时快照解耦", page="09")
    add_picture_fit(slide, ASSET_DIR / "thread_communication_diagram.png", Inches(0.75), Inches(1.35), Inches(7.7), Inches(5.7))
    add_card(slide, 8.75, 1.55, 3.75, 1.2, "采集线程", ["每路输入一个线程，读取摄像头、视频、图片或 RTSP。"], BLUE, 17, 13)
    add_card(slide, 8.75, 3.0, 3.75, 1.2, "推理线程", ["等待 Event，批量取出当前最新帧，串行调用模型。"], GREEN, 17, 13)
    add_card(slide, 8.75, 4.45, 3.75, 1.2, "HTTP 线程", ["ThreadingHTTPServer 处理 /api/events 和 /api/search。"], PURPLE, 17, 13)
    add_card(slide, 8.75, 5.9, 3.75, 0.9, "可靠性", ["最近 100 条事件、20 条错误用于页面展示和排障。"], AMBER, 17, 13)

    # 10. Effect and training section generated by image2.
    slide_full_image(prs, IMAGE2_DIR / "section_03_effect_training_image2.png")

    # 11. Web effect.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Web 实现效果", "用户能直接看到监控、日志、查询结果、路线和车牌裁切图", page="11")
    add_picture_fit(slide, ASSET_DIR / "implementation_effect_diagram.png", Inches(0.65), Inches(1.4), Inches(7.75), Inches(4.9))
    add_card(slide, 8.72, 1.45, 3.85, 1.35, "实时监控", ["4 路摄像头 tile，标注帧以 data URL 返回给前端。"], BLUE, 17, 13)
    add_card(slide, 8.72, 3.05, 3.85, 1.35, "车辆定位", ["输入车牌后返回最后出现摄像头、时间和裁切图。"], GREEN, 17, 13)
    add_card(slide, 8.72, 4.65, 3.85, 1.35, "路线提示", ["前端 SVG 路径高亮入口到目标摄像头点位。"], CYAN, 17, 13)
    add_picture_fit(slide, ASSET_DIR / "sample_test_plate.png", Inches(0.9), Inches(6.05), Inches(2.4), Inches(0.85))
    add_picture_fit(slide, ASSET_DIR / "sample_test_plate2.png", Inches(3.55), Inches(6.05), Inches(2.4), Inches(0.85))

    # 12. Training setup.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "训练配置与数据证据", "训练材料来自 train/，用于证明模型识别链路不是纯演示页面", page="12")
    add_metric(slide, 0.75, 1.55, 2.1, 1.25, "训练轮数", f"{int(final['epoch'])} epoch", BLUE)
    add_metric(slide, 3.15, 1.55, 2.1, 1.25, "输入尺寸", str(args.get("imgsz", "640")), GREEN)
    add_metric(slide, 5.55, 1.55, 2.1, 1.25, "Batch", str(args.get("batch", "256")), CYAN)
    add_metric(slide, 7.95, 1.55, 2.1, 1.25, "模型基座", "YOLO11m", PURPLE)
    add_metric(slide, 10.35, 1.55, 2.1, 1.25, "任务类型", args.get("task", "pose"), AMBER)
    add_picture_fit(slide, ROOT / "train" / "labels.jpg", Inches(0.75), Inches(3.15), Inches(3.55), Inches(3.2))
    add_picture_fit(slide, ROOT / "train" / "val_batch0_pred.jpg", Inches(4.65), Inches(3.15), Inches(3.55), Inches(3.2))
    add_card(slide, 8.55, 3.15, 3.9, 3.2, "训练设计理由", [
        "Pose 任务同时预测车牌框和四角点，服务透视裁切。",
        "640 输入尺寸与板端模型部署尺寸一致。",
        "训练图、PR 曲线和结果 CSV 均保留在 train/ 中可追溯。",
    ], BLUE, 18, 14)

    # 13. Training results.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "训练结果分析", "检测框与关键点指标都达到高水平，支撑后续车牌裁切和识别", page="13")
    add_metric(slide, 0.75, 1.45, 2.25, 1.25, "Box mAP50", f"{final['metrics/mAP50(B)']:.5f}", BLUE)
    add_metric(slide, 3.25, 1.45, 2.25, 1.25, "Box mAP50-95", f"{final['metrics/mAP50-95(B)']:.5f}", CYAN)
    add_metric(slide, 5.75, 1.45, 2.25, 1.25, "Pose mAP50", f"{final['metrics/mAP50(P)']:.5f}", GREEN)
    add_metric(slide, 8.25, 1.45, 2.25, 1.25, "Pose mAP50-95", f"{final['metrics/mAP50-95(P)']:.5f}", PURPLE)
    add_metric(slide, 10.75, 1.45, 1.75, 1.25, "耗时", f"{final['time']/3600:.1f}h", AMBER)
    add_picture_fit(slide, ASSET_DIR / "training_metric_dashboard.png", Inches(0.7), Inches(3.05), Inches(5.75), Inches(3.65))
    add_picture_fit(slide, ASSET_DIR / "training_loss_dashboard.png", Inches(6.75), Inches(3.05), Inches(5.75), Inches(3.65))

    # 14. Deployment section generated by image2.
    slide_full_image(prs, IMAGE2_DIR / "section_04_deployment_outlook_image2.png")

    # 15. Deployment plan.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "部署与运行方式", "PC 用于开发联调，RDK X5 用于现场边缘部署，同一套 Web/SQLite 逻辑复用", page="15")
    add_picture_fit(slide, ROOT / "docs" / "tutorial-assets" / "pc-rdk-deployment.png", Inches(0.75), Inches(1.5), Inches(6.6), Inches(4.9))
    add_card(slide, 7.7, 1.5, 4.8, 1.35, "PC 开发验证", [
        "python web_app.py --backend pc",
        "YOLO .pt + PaddleOCR，适合演示和调参。",
    ], BLUE, 18, 13)
    add_card(slide, 7.7, 3.1, 4.8, 1.35, "RDK 板端部署", [
        "python3 web_app.py --backend bpu",
        "YOLO .bin + LPRNet .bin，调用 BPU runtime。",
    ], GREEN, 18, 13)
    add_card(slide, 7.7, 4.7, 4.8, 1.35, "数据库浏览模式", [
        "--backend none 可打开页面查看历史记录和演示样例。",
    ], AMBER, 18, 13)

    # 16. Highlights.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "项目亮点", "不是单个模型 Demo，而是完整的边缘 AI 应用闭环", page="16")
    highlights = [
        ("双后端架构", ["PC 快速验证，RDK BPU 现场部署，统一上层接口。"], BLUE),
        ("关键点裁切", ["YOLO Pose 输出四角点，透视裁切提升车牌识别可用性。"], GREEN),
        ("实时线程模型", ["最新帧覆盖策略保证低延迟，避免视频队列积压。"], CYAN),
        ("可视化闭环", ["Web 页面同时展示监控、日志、查询、路线和裁切图。"], PURPLE),
        ("训练证据充分", ["400 epoch 训练结果和多类图表保存在 train/ 中。"], AMBER),
        ("工程可维护", ["采集、推理、存储、接口和页面职责分离。"], BLUE),
    ]
    for idx, (title, body, accent) in enumerate(highlights):
        row = idx // 3
        col = idx % 3
        add_card(slide, 0.75 + col * 4.15, 1.55 + row * 2.35, 3.75, 1.85, title, body, accent, 20, 15)

    # 17. Risks and roadmap.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "风险与改进方向", "主动暴露工程风险，并给出后续迭代路径", page="17")
    risks = [
        ("低照度 / 反光", "保存裁切图，分析曝光和角点稳定性，补充夜间样本"),
        ("RTSP 网络抖动", "保持最新帧策略，增加断流重连和状态告警"),
        ("OCR 误识别", "引入置信度、模糊匹配、人工核验和多帧投票"),
        ("数据库增长", "增加归档、去重、保留周期和统计面板"),
        ("部署适配", "对齐 RDK runtime、模型输入尺寸和字符表版本"),
    ]
    y = 1.45
    for i, (risk, fix) in enumerate(risks):
        add_panel(slide, 0.85, y, 11.7, 0.85, accent=BLUE if i % 2 == 0 else GREEN)
        add_text(slide, risk, 1.05, y + 0.16, 2.55, 0.32, size=17, bold=True, color=DEEP_BLUE)
        add_text(slide, fix, 3.75, y + 0.16, 8.4, 0.36, size=15, color=TEXT)
        y += 1.0

    # 18. Closing.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "应用展望", "把车牌识别结果变成可查询、可导航、可运维的地库寻车能力", page="18")
    add_text(slide, "从识别一张车牌，到构建一套可落地的车库车辆定位系统", 0.9, 1.55, 11.5, 0.65, size=30, bold=True, color=DEEP_BLUE, align=PP_ALIGN.CENTER)
    roadmap = [
        ("真实车库地图", "接入楼层、区域和车位拓扑"),
        ("移动端寻车", "面向车主提供快速查询入口"),
        ("跨镜轨迹融合", "从最后位置扩展为完整运动轨迹"),
        ("长期运维告警", "识别异常、摄像头离线和数据归档"),
    ]
    for i, (title, body) in enumerate(roadmap):
        add_flow_step(slide, 1.0 + i * 3.05, 3.0, 2.35, 1.45, title, body, [BLUE, GREEN, CYAN, AMBER][i])
        if i < 3:
            arrow(slide, 3.35 + i * 3.05, 3.72, 3.9 + i * 3.05, 3.72, CYAN)
    add_text(slide, "谢谢观看", 4.75, 5.6, 3.8, 0.55, size=34, bold=True, color=DEEP_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, "Garage Vehicle Locator · Edge AI · Web Visualization", 3.2, 6.25, 6.9, 0.35, size=16, color=MUTED, align=PP_ALIGN.CENTER)

    prs.save(OUT_PATH)
    print(OUT_PATH)


if __name__ == "__main__":
    build_ppt()
